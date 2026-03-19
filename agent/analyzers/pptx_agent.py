import logging
import os
import re
from datetime import date
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from agent.exceptions import PPTXRenderError


def build_output_path(company: str, topic: str, date_str: str, output_dir: str = "output") -> str:
    """Build a stable, human-readable PPTX path.

    Same company + topic + date always returns the same path, enabling in-place revision
    without creating duplicate files.
    """
    def slugify(s: str) -> str:
        s = s.lower().strip()
        s = re.sub(r"[^\w\s-]", "", s)
        s = re.sub(r"[\s_]+", "-", s)
        return s[:40].strip("-")

    filename = f"{slugify(company)}-{slugify(topic)}-{date_str}.pptx"
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    return str(Path(output_dir) / filename)

logger = logging.getLogger(__name__)

# BCG color palette
NAVY = RGBColor(0x0C, 0x20, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0x7D, 0xD3, 0xFC)
ACCENT_BLUE = RGBColor(0x0E, 0xA5, 0xE9)
FOOTER_GRAY = RGBColor(0x94, 0xA3, 0xB8)
BODY_BG = RGBColor(0xFF, 0xFF, 0xFF)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Header bar
HEADER_H = Inches(1.3)

# Body area starts below header
BODY_TOP = Inches(1.45)
BODY_H = Inches(5.6)
BODY_L = Inches(0.5)
BODY_W = Inches(12.33)

# Footer
FOOTER_TOP = Inches(7.1)
FOOTER_H = Inches(0.35)


def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _add_header(slide: Any, title: str, category: str = "") -> None:
    """Add navy header bar with category label and title."""
    # Navy background rectangle
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, HEADER_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Category label (light blue, small caps)
    if category:
        cat_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(10), Inches(0.35))
        tf = cat_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = category.upper()
        run.font.color.rgb = LIGHT_BLUE
        run.font.size = Pt(9)
        run.font.bold = False

    # Title (white, bold, 28pt)
    title_top = Inches(0.42) if category else Inches(0.35)
    title_box = slide.shapes.add_textbox(Inches(0.4), title_top, Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.color.rgb = WHITE
    run.font.size = Pt(26)
    run.font.bold = True


def _add_footer(slide: Any, slide_num: int, total: int = 5) -> None:
    """Add footer with date and slide number."""
    today = date.today().strftime("%B %Y")
    footer_text = f"Competitive Intelligence Report  ·  {today}  ·  {slide_num}/{total}"
    box = slide.shapes.add_textbox(Inches(0.4), FOOTER_TOP, SLIDE_W - Inches(0.8), FOOTER_H)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = footer_text
    run.font.size = Pt(8)
    run.font.color.rgb = FOOTER_GRAY


def _add_body_text(slide: Any, lines: List[str], top: Optional[Any] = None, font_size: int = 12) -> None:
    """Add body text block with bullet lines."""
    t = top if top is not None else BODY_TOP
    box = slide.shapes.add_textbox(BODY_L, t, BODY_W, BODY_H)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = _rgb(0x1E, 0x29, 0x3B)


def _add_section_header(slide: Any, text: str, top: Any) -> Any:
    """Add a colored section header. Returns next top position."""
    box = slide.shapes.add_textbox(BODY_L, top, BODY_W, Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = ACCENT_BLUE
    return top + Inches(0.38)


_PLACEHOLDER = "Insufficient data for this section — expand research scope."


class PPTXAgent:
    """Generates a 5-slide BCG-style PowerPoint from SynthesisAgent output."""

    async def render(self, synthesis: Dict[str, Any], run_id: str) -> str:
        """
        Render a 5-slide BCG deck from synthesis data.

        Output path is derived from company + topic + date so the same research
        always writes to the same file. Revisions overwrite in-place — no duplicates.

        Args:
            synthesis: Output from SynthesisAgent.run()
            run_id: Stored in footer / metadata only

        Returns:
            Absolute path to the generated .pptx file
        """
        try:
            return self._build(synthesis, run_id)
        except Exception as e:
            logger.error("PPTXAgent render error: %s", e)
            raise PPTXRenderError(f"Failed to render deck: {e}") from e

    def _build(self, synthesis: Dict[str, Any], run_id: str) -> str:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        blank_layout = prs.slide_layouts[6]  # Completely blank

        # Slide 1: Executive Summary
        self._slide_executive_summary(prs, blank_layout, synthesis)

        # Slide 2: Market Landscape
        self._slide_market_landscape(prs, blank_layout, synthesis)

        # Slide 3: Competitive Analysis
        self._slide_competitive_analysis(prs, blank_layout, synthesis)

        # Slide 4: Strategic Implications
        self._slide_strategic_implications(prs, blank_layout, synthesis)

        # Slide 5: Recommendations & Outlook
        self._slide_recommendations(prs, blank_layout, synthesis)

        # Named output — stable path enables in-place revision
        company = synthesis.get("company_name") or synthesis.get("target") or "research"
        topic = synthesis.get("topic") or synthesis.get("scope") or "intelligence-brief"
        date_str = date.today().isoformat()
        out_path_str = build_output_path(company, topic, date_str)
        out_path = Path(out_path_str)

        prs.save(str(out_path))
        logger.info("PPTXAgent: saved deck to %s", out_path)
        return str(out_path.resolve())

    def _slide_executive_summary(self, prs: Presentation, layout: Any, s: Dict[str, Any]) -> None:
        slide = prs.slides.add_slide(layout)
        _add_header(slide, "Executive Summary", "COMPETITIVE INTELLIGENCE")
        _add_footer(slide, 1)

        summary = s.get("executive_summary") or _PLACEHOLDER
        recs = s.get("recommendations") or []

        lines: List[str] = [f"• {summary}"]
        if recs:
            lines.append("")
            lines.append("Key Recommendations:")
            for rec in recs[:3]:
                lines.append(f"  › {rec}")

        _add_body_text(slide, lines)

    def _slide_market_landscape(self, prs: Presentation, layout: Any, s: Dict[str, Any]) -> None:
        slide = prs.slides.add_slide(layout)
        _add_header(slide, "Market Landscape", "MARKET INTELLIGENCE")
        _add_footer(slide, 2)

        ml = s.get("market_landscape") or {}
        size_growth = ml.get("size_and_growth") or _PLACEHOLDER
        players: List[Dict[str, Any]] = ml.get("key_players") or []
        trends: List[str] = ml.get("trends") or []

        lines: List[str] = [f"• {size_growth}", ""]

        if players:
            lines.append("Key Players:")
            for p in players[:5]:
                name = p.get("name", "")
                position = p.get("position", "")
                signal = p.get("signal", "")
                lines.append(f"  › {name}: {position} — {signal}")
            lines.append("")

        if trends:
            lines.append("Market Trends:")
            for t in trends[:5]:
                lines.append(f"  › {t}")

        _add_body_text(slide, lines)

    def _slide_competitive_analysis(self, prs: Presentation, layout: Any, s: Dict[str, Any]) -> None:
        slide = prs.slides.add_slide(layout)
        _add_header(slide, "Competitive Analysis", "COMPETITIVE POSITIONING")
        _add_footer(slide, 3)

        ca = s.get("competitive_analysis") or {}
        table: List[Dict[str, Any]] = ca.get("comparison_table") or []
        winners: List[str] = ca.get("winner_signals") or []
        disruptions: List[str] = ca.get("disruption_risks") or []

        lines: List[str] = []
        if table:
            lines.append("Competitive Comparison:")
            for row in table[:6]:
                dim = row.get("dimension", "")
                findings = row.get("findings", "")
                lines.append(f"  {dim}: {findings}")
            lines.append("")
        else:
            lines.append(_PLACEHOLDER)
            lines.append("")

        if winners:
            lines.append("Winner Signals:")
            for w in winners[:4]:
                lines.append(f"  › {w}")
            lines.append("")

        if disruptions:
            lines.append("Disruption Risks:")
            for d in disruptions[:4]:
                lines.append(f"  › {d}")

        _add_body_text(slide, lines)

    def _slide_strategic_implications(self, prs: Presentation, layout: Any, s: Dict[str, Any]) -> None:
        slide = prs.slides.add_slide(layout)
        _add_header(slide, "Strategic Implications", "STRATEGIC ANALYSIS")
        _add_footer(slide, 4)

        si = s.get("strategic_implications") or {}
        opps: List[str] = si.get("opportunities") or []
        risks: List[str] = si.get("risks") or []
        watch: List[str] = si.get("watch_list") or []

        lines: List[str] = []
        if opps:
            lines.append("Opportunities:")
            for o in opps[:4]:
                lines.append(f"  › {o}")
            lines.append("")
        else:
            lines.append("Opportunities: " + _PLACEHOLDER)
            lines.append("")

        if risks:
            lines.append("Risks:")
            for r in risks[:4]:
                lines.append(f"  › {r}")
            lines.append("")
        else:
            lines.append("Risks: " + _PLACEHOLDER)
            lines.append("")

        if watch:
            lines.append("Watch List:")
            for w in watch[:4]:
                lines.append(f"  › {w}")

        _add_body_text(slide, lines)

    def _slide_recommendations(self, prs: Presentation, layout: Any, s: Dict[str, Any]) -> None:
        slide = prs.slides.add_slide(layout)
        _add_header(slide, "Recommendations & Outlook", "STRATEGIC DIRECTION")
        _add_footer(slide, 5)

        recs: List[str] = s.get("recommendations") or []
        outlook = s.get("outlook") or _PLACEHOLDER

        lines: List[str] = []
        if recs:
            lines.append("Recommendations:")
            for i, r in enumerate(recs, 1):
                lines.append(f"  {i}. {r}")
            lines.append("")
        else:
            lines.append("Recommendations: " + _PLACEHOLDER)
            lines.append("")

        lines.append("Outlook:")
        lines.append(f"  {outlook}")

        _add_body_text(slide, lines)
